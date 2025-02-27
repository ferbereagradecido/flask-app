from flask import Flask, request, jsonify, render_template
import openai
import os
import sqlite3
import re
from werkzeug.utils import secure_filename
import PyPDF2
from docx import Document
from pptx import Presentation

# Initialize Flask app
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads/'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Initialize SQLite database
DATABASE = 'responses.db'

def init_db():
    with sqlite3.connect(DATABASE) as conn:
        cursor = conn.cursor()
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS responses (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                question TEXT,
                user_answer TEXT,
                feedback TEXT
            )
        ''')
        conn.commit()
init_db()

# Initialize OpenAI client
openai.api_key = os.getenv("OPENAI_API_KEY")

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/process', methods=['POST'])
def process():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded."})
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file."})
    
    num_questions = request.form.get('numQuestions')
    if not num_questions:
        return jsonify({"error": "Number of questions not provided."})
    num_questions = int(num_questions)
    
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)
    
    file_ext = filename.rsplit('.', 1)[1].lower()
    document_text = ""
    
    if file_ext == 'pdf':
        document_text = extract_text_from_pdf(filepath)
    elif file_ext == 'docx':
        document_text = extract_text_from_docx(filepath)
    elif file_ext in ['ppt', 'pptx']:
        document_text = extract_text_from_pptx(filepath)
    else:
        return jsonify({"error": "Unsupported file format."})
    
    summary = generate_summary_outline(document_text)
    questions = generate_questions(document_text, num_questions)
    
    return jsonify({
        "summary": summary,
        "questions": questions,
        "num_questions": num_questions
    })

@app.route('/evaluate_answer', methods=['POST'])
def evaluate_answer():
    data = request.json
    question = data.get("question")
    user_answer = data.get("userAnswer")
    document_text = data.get("documentText")
    
    if not question or not user_answer or not document_text:
        return jsonify({"error": "Missing data."})
    
    response = openai.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a strict evaluator who provides feedback based on the given document."},
            {"role": "user", "content": f"Document:\n{document_text[:4000]}\n\nQuestion:\n{question}\n\nUser's answer:\n{user_answer}\n\nEvaluate the answer strictly based on the document. Provide a score from 1 to 10, what is correct, what is missing, and how the user can improve."}
        ]
    )
    feedback = response.choices[0].message.content
    
    # Save response to the database
    with sqlite3.connect(DATABASE) as conn:
        cursor = conn.cursor()
        cursor.execute("INSERT INTO responses (question, user_answer, feedback) VALUES (?, ?, ?)", (question, user_answer, feedback))
        conn.commit()
    
    return jsonify({"feedback": feedback})

@app.route('/responses', methods=['GET'])
def get_responses():
    with sqlite3.connect(DATABASE) as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT question, user_answer, feedback FROM responses")
        responses = cursor.fetchall()
    return jsonify({"responses": responses})

# Functions to extract text from different file types
def extract_text_from_pdf(pdf_path):
    text = ""
    with open(pdf_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Generate a summary outline using OpenAI
def generate_summary_outline(document_text):
    response = openai.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "user", "content": f"Summarize this text:\n{document_text[:4000]}"}
        ]
    )
    return response.choices[0].message.content

# Generate questions from the document
def generate_questions(document_text, num_questions=1):
    response = openai.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "user", "content": f"Generate {num_questions} questions from this document:\n{document_text[:4000]}"}
        ]
    )
    return response.choices[0].message.content.split("\n")

if __name__ == '__main__':
    app.run(debug=True)
